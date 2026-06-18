"""PPTX Generator — 마크다운 보고서 → PPT 슬라이드 (Phase3, 2026-06-09).

python-pptx 사용(한국어 네이티브). 렌더 단계 tool(ToolCategory.RENDERING).
report_markdown 을 받아: 첫 '# ' = 타이틀 슬라이드, 각 '## ' = 섹션 슬라이드(이후 줄=불릿).
chart_image_paths 는 선택 — 있으면 섹션 뒤에 차트 슬라이드로 첨부, 없어도 동작.
(2026-06-12: 실구현 — 구 chart_to_slide stub 의 "차트→슬라이드 배치" 책임을 여기로 흡수.
 별도 tool 이던 chart_to_slide 는 산출(chart_slides) 소비자 0 으로 폐기, 헌법 R6.)
빈 입력(report_markdown 부재)은 data_insufficient.
"""
from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Any

from app.core.logging import get_logger
from app.dream_agent.models import ExecutionContext
from app.dream_agent.tools.base_tool import BaseTool
from app.dream_agent.tools.shared.helpers import find_in_previous

logger = get_logger(__name__)

# pptx_generator.py → rendering(0) tools(1) dream_agent(2) app(3) backend(4) repo(5)
_REPO_ROOT = Path(__file__).resolve().parents[5]


def _strip_bullet(s: str) -> str:
    return re.sub(r"^\s*[-*]\s+", "", s.strip())


_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")


def _strip_inline_md(s: str) -> str:
    """타이틀용 — 인라인 마크다운 마커(굵게/코드) 제거, 텍스트만 (날것 ** 방지)."""
    s = _BOLD_RE.sub(r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    return s.strip()


def _segments(text: str) -> list[tuple[str, bool]]:
    """인라인 마크다운 → (텍스트, 굵게여부) 세그먼트.

    '`a **b** c`' → [('a ', False), ('b', True), (' c', False)].
    **x** 만 굵게로 분리하고, 비굵게 구간의 코드 마커(`)는 텍스트로 정리 → 슬라이드에 날것 ** 0.
    """
    out: list[tuple[str, bool]] = []
    pos = 0
    for m in _BOLD_RE.finditer(text):
        if m.start() > pos:
            out.append((text[pos : m.start()], False))
        out.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], False))
    cleaned = [
        (re.sub(r"`([^`]+)`", r"\1", seg) if not bold else seg, bold)
        for seg, bold in out
    ]
    return [(s, b) for s, b in cleaned if s] or [(text, False)]


def _set_runs(p, text: str) -> None:
    """문단(p)에 인라인 마크다운 반영 — **굵게**=볼드 run, 나머지=일반 run."""
    for seg, bold in _segments(text):
        run = p.add_run()
        run.text = seg
        if bold:
            run.font.bold = True


def _render_pptx(
    markdown_text: str,
    out_path: Path,
    charts: list[tuple[str, str]] | None = None,
) -> int:
    """마크다운 → 슬라이드. 슬라이드 수 반환. (# 타이틀 + ## 섹션별 불릿 + 차트 슬라이드)

    charts: (제목, 이미지경로) — 존재하는 파일만 슬라이드로. 없으면 텍스트 전용(기존 동작).
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    lines = markdown_text.splitlines()

    title = next((l[2:].strip() for l in lines if l.startswith("# ")), "분석 보고서")
    ts = prs.slides.add_slide(prs.slide_layouts[0])  # 타이틀 레이아웃
    ts.shapes.title.text = _strip_inline_md(title)

    # ## 섹션별 (title, [bullets]) 수집
    sections: list[tuple[str, list[str]]] = []
    cur: tuple[str, list[str]] | None = None
    for l in lines:
        if l.startswith("## "):
            cur = (l[3:].strip(), [])
            sections.append(cur)
        elif l.startswith("# "):
            continue
        elif l.strip():
            if cur is None:
                cur = ("개요", [])
                sections.append(cur)
            cur[1].append(_strip_bullet(l))

    for sec_title, bullets in sections:
        sl = prs.slides.add_slide(prs.slide_layouts[1])  # 제목+내용
        sl.shapes.title.text = _strip_inline_md(sec_title)
        tf = sl.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets or ["—"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _set_runs(p, b)

    # 차트 슬라이드 — 실재 파일만 (없는 경로를 빈 슬라이드로 꾸미지 않는다, I1)
    chart_n = 0
    for title, img in charts or []:
        if not Path(img).exists():
            logger.warning("pptx chart 첨부 생략 — 파일 부재", path=img)
            continue
        sl = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        sl.shapes.title.text = title
        sl.shapes.add_picture(img, Inches(0.6), Inches(1.5), width=Inches(8.8))
        chart_n += 1

    prs.save(str(out_path))
    return 1 + len(sections) + chart_n


class PptxGenerator(BaseTool):
    """report_markdown → PPT 슬라이드 파일. 렌더 단계(ToolCategory.RENDERING)."""

    async def execute(
        self,
        params: dict[str, Any],
        context: ExecutionContext,
    ) -> dict[str, Any]:
        previous = context.previous_results or {}
        markdown = (
            find_in_previous(previous, "report_markdown")
            or params.get("report_markdown")
            or ""
        )
        if not markdown:
            logger.info("pptx_generator skipped — report_markdown 부재(data_insufficient)",
                        session_id=context.session_id)
            return {"pptx_file_path": None, "reason": "data_insufficient",
                    "detail": "report_markdown 0건/부재"}

        out_dir = Path(
            params.get("output_dir")
            or (_REPO_ROOT / "data" / (context.client_id or "default") / "outputs")
        )
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        # 차트(선택) — chart_generator 산출. charts 메타(제목)와 경로를 index 로 짝지음.
        chart_paths = (
            find_in_previous(previous, "chart_image_paths")
            or params.get("chart_image_paths") or []
        )
        charts_meta = find_in_previous(previous, "charts") or []
        titles = [
            (charts_meta[i].get("title") if i < len(charts_meta) and isinstance(charts_meta[i], dict) else None)
            or f"차트 {i + 1}"
            for i in range(len(chart_paths))
        ]
        charts = list(zip(titles, [str(p) for p in chart_paths]))

        slides = _render_pptx(markdown, out_path, charts)
        logger.info("pptx_generator completed", path=str(out_path), slides=slides,
                    charts=len(charts))
        return {"pptx_file_path": str(out_path), "slide_count": slides}
