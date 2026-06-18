"""Phase 3 — pptx 인라인 마크다운 정리 (다운로드 파일 품질) (2026-06-09).

문제: report_markdown 의 인라인 '**굵게**' 가 슬라이드 텍스트에 날것 `**` 로 노출됨.
(채팅 시각화는 react-markdown 으로 처리되지만, 다운로드한 .pptx 파일은 별도 렌더.)
PDF 는 이미 **→<b> 변환됨(pdf_renderer). PPTX 만 미처리였음.

M-1 _segments: **x** → 굵게 분리, 코드 마커 정리
M-2 렌더된 .pptx 의 어떤 run 텍스트에도 `**` 없음 + 굵게 run 존재
M-3 타이틀/섹션 제목의 인라인 마커도 제거
"""
from __future__ import annotations

from app.dream_agent.tools.rendering.pptx_generator import _render_pptx, _segments


# ── M-1: 세그먼트 분리 ──

def test_m1_segments_splits_bold():
    assert _segments("a **b** c") == [("a ", False), ("b", True), (" c", False)]


def test_m1_segments_plain_passthrough():
    assert _segments("마커 없음") == [("마커 없음", False)]


def test_m1_segments_strips_code_marker():
    # 코드 마커는 텍스트만 남김 (날것 ` 제거)
    assert _segments("값 `metric` 확인") == [("값 metric 확인", False)]


# ── M-2: 렌더 파일에 날것 ** 없음 + 굵게 run ──

def test_m2_pptx_no_raw_asterisks_and_has_bold(tmp_path):
    from pptx import Presentation

    md = "# 4월 보고서\n## 핵심\n- **긍정 58.3%**로 높음\n- 일반 항목"
    out = tmp_path / "t.pptx"
    _render_pptx(md, out)

    prs = Presentation(str(out))
    all_text: list[str] = []
    bold_text: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    all_text.append(run.text)
                    if run.font.bold:
                        bold_text.append(run.text)

    joined = "".join(all_text)
    assert "**" not in joined, f"날것 ** 노출: {joined!r}"
    assert "긍정 58.3%" in bold_text, "굵게 구간이 볼드 run 으로 분리되어야 함"


# ── M-3: 제목의 인라인 마커도 제거 ──

def test_m3_title_strips_inline_marker(tmp_path):
    from pptx import Presentation

    md = "# **4월** 보고서\n## 핵심 `요약`\n- 항목"
    out = tmp_path / "t.pptx"
    _render_pptx(md, out)

    prs = Presentation(str(out))
    titles = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame and shape == slide.shapes.title
    ]
    joined = " ".join(titles)
    assert "**" not in joined and "`" not in joined, f"제목 날것 마커 노출: {joined!r}"
    assert "4월" in joined and "핵심" in joined
